"""
Generates:
  docs/presentation.pptx  — polished light-theme 10-slide exam deck
  docs/talking_paper.pdf  — full speaker notes + appendix

Exam format (SYS Spring 2026, Henrik Kuhl):
  Presentation  ~10 min   (~60 sec/slide × 10 slides)
  Discussion    ~15 min   (synopsis + curriculum topics)
  Total          30 min
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak,
    HRFlowable, KeepTogether, Table, TableStyle,
)
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

OUT  = os.path.dirname(os.path.abspath(__file__))

# ── Palette ──────────────────────────────────────────────────────────────────
BG    = RGBColor(0xF4, 0xF7, 0xFB)
HDR   = RGBColor(0x10, 0x20, 0x38)
ACNT  = RGBColor(0x00, 0x88, 0xAA)
TEXT  = RGBColor(0x10, 0x20, 0x38)
LGRAY = RGBColor(0x58, 0x68, 0x7C)
MGRAY = RGBColor(0xC4, 0xD2, 0xE0)
GREEN = RGBColor(0x12, 0x70, 0x38)
RED   = RGBColor(0xB4, 0x16, 0x16)
ORANGE= RGBColor(0x9C, 0x4E, 0x04)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DBOX  = RGBColor(0xE2, 0xEC, 0xF8)
CBOX  = RGBColor(0x14, 0x22, 0x36)

C_API = RGBColor(0x00, 0x54, 0x92)
C_STK = RGBColor(0x00, 0x66, 0x66)
C_PAY = RGBColor(0x4E, 0x20, 0x82)
C_FAN = RGBColor(0x1E, 0x5E, 0x26)
C_EVT = RGBColor(0x1A, 0x3A, 0x6A)
C_INF = RGBColor(0x66, 0x3C, 0x06)

LT_RED = RGBColor(0xFC, 0xEB, 0xEB)
LT_GRN = RGBColor(0xE6, 0xF7, 0xEB)
LT_BLU = RGBColor(0xE2, 0xEC, 0xF8)
LT_ORG = RGBColor(0xFC, 0xF2, 0xE2)

FONT = "Calibri"


# ── Primitives ────────────────────────────────────────────────────────────────

def _bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _rect(slide, l, t, w, h, fill=DBOX, stroke=None, sw=Pt(1)):
    """Sharp-corner filled rectangle."""
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if stroke:
        sh.line.color.rgb = stroke
        sh.line.width = sw
    else:
        sh.line.fill.background()
    return sh


def _label(slide, text, l, t, w, h,
           fill=C_EVT, stroke=None, sw=Pt(1.2),
           size=11, bold=True, color=WHITE,
           align=PP_ALIGN.CENTER, wrap=True):
    """Sharp-corner box with centred label text."""
    sh = _rect(slide, l, t, w, h, fill=fill, stroke=stroke, sw=sw)
    tf = sh.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    return sh


def _txt(slide, text, l, t, w, h,
         size=13, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    r.font.color.rgb = color
    return txb


def _arrow(slide, x1, y1, x2, y2, color=ACNT, w=Pt(1.5)):
    """Connector with arrowhead at (x2, y2)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = w
    sp   = conn._element
    spPr = sp.find(qn('p:spPr'))
    ln   = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    for te in ln.findall(qn('a:tailEnd')):
        ln.remove(te)
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow')
    te.set('w', 'med')
    te.set('len', 'med')
    return conn


def _line(slide, x1, y1, x2, y2, color=MGRAY, w=Pt(1.2)):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = w
    return conn


def _header(slide, title, sub=None):
    _rect(slide, 0, 0, 13.33, 1.08, fill=HDR)
    _rect(slide, 0, 1.00, 13.33, 0.08, fill=ACNT)
    _txt(slide, title, 0.32, 0.08, 12.70, 0.70, size=26, bold=True, color=WHITE)
    if sub:
        _txt(slide, sub, 0.36, 0.74, 12.70, 0.24, size=11, color=MGRAY, italic=True)


def _footer(slide, n, total=12):
    pw = 13.33 * n / total
    _rect(slide, 0, 7.46, 13.33, 0.04, fill=MGRAY)
    _rect(slide, 0, 7.46, pw,    0.04, fill=ACNT)
    _txt(slide, f"{n} / {total}", 12.35, 7.28, 0.82, 0.18,
         size=9, color=LGRAY, align=PP_ALIGN.RIGHT)


def _code_box(slide, lines, l, t, w, h, size=9.5):
    KW  = {'using','public','private','class','namespace','async','await',
           'return','var','new','if','else','try','catch','throw','void',
           'Task','bool','true','false','null','int','string','Guid',
           'DateTime','static','readonly','override','base','record'}
    KW_C  = RGBColor(0x56, 0x9C, 0xD6)
    STR_C = RGBColor(0xCE, 0x91, 0x78)
    CMT_C = RGBColor(0x6A, 0x99, 0x55)
    DEF_C = RGBColor(0xD4, 0xD4, 0xD4)
    TYP_C = RGBColor(0x4E, 0xC9, 0xB0)

    _rect(slide, l, t, w, h, fill=CBOX, stroke=ACNT, sw=Pt(0.75))
    txb = slide.shapes.add_textbox(
        Inches(l+0.13), Inches(t+0.11), Inches(w-0.26), Inches(h-0.22))
    tf = txb.text_frame
    tf.word_wrap = False
    tf.auto_size = None

    def _run(para, text, c):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = "Courier New"
        r.font.color.rgb = c

    def _tok(line):
        s = line.lstrip()
        if s.startswith('//') or s.startswith('#'):
            return [(line, CMT_C)]
        ci = line.find('//')
        if ci > 0 and line[:ci].count('"') % 2 == 0:
            return _seg(line[:ci]) + [(line[ci:], CMT_C)]
        return _seg(line)

    def _seg(seg):
        out = []
        for part in re.split(r'("(?:[^"\\]|\\.)*")', seg):
            if part.startswith('"') and part.endswith('"') and len(part) >= 2:
                out.append((part, STR_C))
            else:
                for wp in re.split(r'(\b\w+\b)', part):
                    if wp in KW:
                        out.append((wp, KW_C))
                    elif re.match(r'^\w+$', wp) and wp and wp[0].isupper():
                        out.append((wp, TYP_C))
                    else:
                        out.append((wp, DEF_C))
        return out

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        if not line.strip():
            _run(p, ' ', DEF_C)
            continue
        for text, c in _tok(line):
            if text:
                _run(p, text, c)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _rect(slide, 0, 0,    13.33, 1.70, fill=HDR)
    _rect(slide, 0, 1.62, 13.33, 0.08, fill=ACNT)
    _rect(slide, 0, 6.48, 13.33, 1.02, fill=HDR)
    _rect(slide, 0, 6.41, 13.33, 0.07, fill=ACNT)
    _rect(slide, 0,     0, 0.06, 7.5, fill=ACNT)
    _rect(slide, 13.27, 0, 0.06, 7.5, fill=ACNT)

    _txt(slide, "Event-Driven Order\nFulfillment System",
         1.0, 1.80, 11.33, 1.80,
         size=46, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    _txt(slide, "System Integration  ·  Spring 2026",
         1.0, 3.68, 11.33, 0.52,
         size=20, color=ACNT, align=PP_ALIGN.CENTER, bold=True)

    _rect(slide, 3.5, 4.34, 6.33, 0.04, fill=MGRAY)

    # 6 equal-width tech tags centred in the slide
    tags = [".NET 10", "MassTransit 8", "RabbitMQ",
            "PostgreSQL", "Docker Compose", "xUnit / Moq / WireMock"]
    TW, GAP = 1.88, 0.22          # 6×1.88 + 5×0.22 = 12.38"  start x = 0.475"
    tx = (13.33 - (len(tags)*TW + (len(tags)-1)*GAP)) / 2
    for tag in tags:
        _label(slide, tag, tx, 4.52, TW, 0.44,
               fill=LT_BLU, stroke=ACNT, color=C_EVT, size=11, bold=True, sw=Pt(1))
        tx += TW + GAP

    _txt(slide, "Asyno  ·  PBSW  ·  2026-06-15",
         1.0, 6.58, 11.33, 0.34,
         size=12, color=MGRAY, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "[~20 sec] "
        "Good morning. My name is Marco and I will present my synopsis project: "
        "an Event-Driven Order Fulfillment System in .NET 10. "
        "I will walk through the architecture, four design patterns, and the testing strategy "
        "in about ten minutes."
    )


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Why Event-Driven?",
            "Tight synchronous chains break under load — async events decouple services")
    _footer(slide, 2, total=10)

    # LEFT — synchronous
    _rect(slide, 0.30, 1.20, 5.98, 6.04, fill=LT_RED, stroke=RED, sw=Pt(1.2))
    _txt(slide, "Synchronous Chain",
         0.52, 1.28, 5.50, 0.34, size=14, bold=True, color=RED)

    svcs = ["Order API", "Stock Service", "Payment Service",
            "Warehouse", "Shipping", "Notification"]
    sy = 1.74
    for i, svc in enumerate(svcs):
        _label(slide, svc, 0.62, sy, 5.18, 0.54,
               fill=RGBColor(0xB0,0x1A,0x1A), stroke=None, size=12, sw=Pt(0))
        if i < len(svcs) - 1:
            _arrow(slide, 3.21, sy+0.54, 3.21, sy+0.68, color=RED, w=Pt(1.2))
        sy += 0.68

    _txt(slide,
         "Response time = sum of ALL latencies\n"
         "One failure cascades to every caller",
         0.52, 6.32, 5.50, 0.72, size=11, color=RED)

    _line(slide, 6.64, 1.22, 6.64, 7.18, color=MGRAY, w=Pt(0.75))

    # RIGHT — event-driven
    _rect(slide, 6.76, 1.20, 6.28, 6.04, fill=LT_GRN, stroke=GREEN, sw=Pt(1.2))
    _txt(slide, "Event-Driven (async)",
         6.98, 1.28, 5.80, 0.34, size=14, bold=True, color=GREEN)

    _label(slide, "Order API", 7.16, 1.74, 5.70, 0.54,
           fill=C_API, stroke=None, size=13, sw=Pt(0))
    _txt(slide, "publishes  OrderPlaced",
         7.40, 2.34, 5.20, 0.26, size=10.5, color=GREEN, italic=True)
    _arrow(slide, 10.01, 2.60, 10.01, 2.82, color=GREEN, w=Pt(1.3))
    _label(slide, "RabbitMQ  (message broker)", 7.16, 2.85, 5.70, 0.54,
           fill=C_INF, stroke=None, size=12, sw=Pt(0))
    _txt(slide, "routes a copy to every subscriber:",
         7.40, 3.45, 5.20, 0.26, size=10.5, color=GREEN, italic=True)
    _arrow(slide, 10.01, 3.71, 10.01, 3.94, color=GREEN, w=Pt(1.3))

    names = ["Stock", "Payment", "Warehouse", "Shipping", "Notif."]
    n   = len(names)
    sw2 = (5.70 - (n-1)*0.12) / n
    BUS = 3.97
    sx0 = 7.16
    bus_r = sx0 + n*(sw2+0.12) - 0.12
    _line(slide, sx0, BUS, bus_r, BUS, color=GREEN, w=Pt(1.8))
    for i, nm in enumerate(names):
        lx = sx0 + i*(sw2+0.12)
        cx = lx + sw2/2
        _arrow(slide, cx, BUS, cx, 4.14, color=GREEN, w=Pt(1.2))
        _label(slide, nm, lx, 4.17, sw2, 0.52,
               fill=C_FAN, stroke=None, size=10, sw=Pt(0))

    _txt(slide,
         "HTTP 201 immediately — processing is async\n"
         "Services scale and fail completely independently",
         6.98, 4.88, 5.80, 0.72, size=11, color=GREEN)

    slide.notes_slide.notes_text_frame.text = (
        "[~60 sec] "
        "Synchronously every service call blocks the next — the customer waits for the "
        "sum of all latencies, and if Shipping goes down the entire chain fails. "
        "Event-driven: the Order API publishes one event and returns HTTP 201 immediately. "
        "RabbitMQ routes a copy to every subscriber in parallel. "
        "Services are decoupled in time and space — one failure no longer cascades."
    )


def slide_03_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "System Architecture",
            "Choreography saga — each service publishes the next event")
    _footer(slide, 3, total=10)

    # Saga chain — 4 boxes, full width centred
    BW, BH, GAP = 2.72, 0.74, 0.46
    lm = (13.33 - (4*BW + 3*GAP)) / 2   # 0.535"

    chain = [
        ("Order Service API", C_API),
        ("Stock Consumer",    C_STK),
        ("Payment Consumer",  C_PAY),
        ("Order Service\n(confirms)", C_API),
    ]
    BOX_T = 1.96
    cx_list = []
    for i, (name, col) in enumerate(chain):
        lx = lm + i*(BW+GAP)
        _label(slide, name, lx, BOX_T, BW, BH, fill=col, stroke=None, size=12, sw=Pt(0))
        cx_list.append(lx + BW/2)

    arr_y = BOX_T + BH/2
    for i in range(3):
        _arrow(slide, cx_list[i]+BW/2, arr_y, cx_list[i+1]-BW/2, arr_y,
               color=ACNT, w=Pt(1.8))

    # Event name badges above arrows
    evts = ["OrderPlaced", "StockReserved", "PaymentSucceeded"]
    for i, evt in enumerate(evts):
        mid = (cx_list[i] + cx_list[i+1]) / 2
        ew  = 1.76
        _label(slide, evt, mid-ew/2, 1.64, ew, 0.28,
               fill=DBOX, stroke=ACNT, color=C_EVT, size=9, bold=True, sw=Pt(0.8))

    # Fan-out downward from rightmost box
    os_cx = cx_list[3]
    _arrow(slide, os_cx, BOX_T+BH, os_cx, 3.52, color=GREEN, w=Pt(1.8))

    ew2 = 2.14
    _label(slide, "OrderConfirmed", os_cx-ew2/2, 3.56, ew2, 0.30,
           fill=LT_GRN, stroke=GREEN, color=GREEN, size=9.5, bold=True, sw=Pt(1))

    # 4 fan-out consumers via horizontal bus
    FW, FH = 2.58, 0.72
    fans = ["Invoice\nConsumer", "Warehouse\nConsumer",
            "Shipping\nConsumer", "Notification\nConsumer"]
    fan_total = 4*FW + 3*0.20
    fl = (13.33 - fan_total) / 2
    fan_cxs = [fl + i*(FW+0.20) + FW/2 for i in range(4)]

    BUS_Y = 3.90
    _line(slide, fan_cxs[0], BUS_Y, fan_cxs[3], BUS_Y, color=GREEN, w=Pt(2))
    _line(slide, os_cx, 3.88, os_cx, BUS_Y, color=GREEN, w=Pt(2))

    for cx, name in zip(fan_cxs, fans):
        _arrow(slide, cx, BUS_Y, cx, 4.12, color=GREEN, w=Pt(1.4))
        _label(slide, name, cx-FW/2, 4.14, FW, FH, fill=C_FAN, stroke=None, size=11.5, sw=Pt(0))

    _rect(slide, 0.30, 5.08, 12.74, 0.44, fill=LT_RED, stroke=RED, sw=Pt(0.9))
    _txt(slide,
         "Failure path:  PaymentFailed -> StockConsumer publishes StockReleased  "
         "+  OrderService publishes OrderCancelled  (compensating transactions)",
         0.50, 5.13, 12.34, 0.36, size=11, color=RED)

    _txt(slide,
         "All events are C# records in the shared Contracts project  "
         "— compile-time type safety across all 7 services",
         0.30, 5.66, 12.74, 0.26, size=10.5, color=LGRAY,
         align=PP_ALIGN.CENTER, italic=True)

    infra = [("RabbitMQ 3", C_INF), ("PostgreSQL 16", C_INF), ("MassTransit 8.3", C_API)]
    ix = 0.38
    for tag, col in infra:
        iw = len(tag)*0.115 + 0.42
        _label(slide, tag, ix, 6.06, iw, 0.38, fill=col, stroke=None, size=10, bold=True, sw=Pt(0))
        ix += iw + 0.20

    slide.notes_slide.notes_text_frame.text = (
        "[~75 sec] "
        "Seven microservices, one shared Contracts library, RabbitMQ as the broker, and "
        "PostgreSQL for state. The horizontal chain is the core saga: OrderService publishes "
        "OrderPlaced, StockConsumer reserves stock, PaymentConsumer charges the card, then "
        "OrderService publishes OrderConfirmed. That single event fans out via a RabbitMQ "
        "fanout exchange to four consumers in parallel — Invoice, Warehouse, Shipping, "
        "Notification — each with its own named queue. "
        "Failure path: if payment fails, compensating events undo earlier steps."
    )


def slide_04_sync_async(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Synchronous vs. Asynchronous Communication",
            "Choosing the right channel is the first design decision")
    _footer(slide, 4, total=10)

    rows = [
        ("Property",       "Synchronous",          "Async / Event-Driven"),
        ("Coupling",       "Tight — direct calls",  "Loose — via broker"),
        ("Response time",  "Sum of all steps",      "Immediate ACK"),
        ("Failure impact", "Cascades to caller",    "Isolated per service"),
        ("Scalability",    "Scale all at once",     "Scale each service"),
        ("Consistency",    "Strong (done now)",     "Eventual (done later)"),
        ("Debugging",      "Simple call stack",     "Distributed log search"),
    ]
    col_x = [0.30, 2.44, 4.78]
    col_w = [2.10, 2.30, 2.30]
    RH    = 0.52

    for ci, (cx, cw, hf) in enumerate(zip(col_x, col_w, [HDR, C_PAY, C_STK])):
        _label(slide, rows[0][ci], cx, 1.20, cw, 0.42,
               fill=hf, stroke=None, size=10.5, bold=True, sw=Pt(0))

    for ri, row in enumerate(rows[1:], 1):
        bg = RGBColor(0xF8,0xF9,0xFF) if ri % 2 == 0 else WHITE
        for ci, (cx, cw) in enumerate(zip(col_x, col_w)):
            if ci == 0:
                _label(slide, row[ci], cx, 1.20+ri*RH, cw, RH-0.04,
                       fill=DBOX, stroke=MGRAY, color=TEXT,
                       size=10.5, bold=True, sw=Pt(0.5))
            else:
                tc = RED if ci == 1 else GREEN
                _label(slide, row[ci], cx, 1.20+ri*RH, cw, RH-0.04,
                       fill=bg, stroke=MGRAY, color=tc,
                       size=10, bold=False, sw=Pt(0.5))

    _rect(slide, 0.30, 5.16, 6.88, 0.48, fill=LT_BLU, stroke=ACNT, sw=Pt(1))
    _txt(slide, "Rule:  choose async when you do NOT need the result right now",
         0.50, 5.24, 6.68, 0.30, size=12, bold=True, color=ACNT)

    for i, (lbl, txt) in enumerate([
        ("Use sync when:",  "Login, balance check — result needed immediately"),
        ("Use async when:", "Email, invoice, label — can happen in background"),
    ]):
        _txt(slide, lbl, 0.42, 5.76+i*0.46, 1.90, 0.28, size=11, bold=True)
        _txt(slide, txt, 2.36, 5.76+i*0.46, 4.80, 0.28, size=11, color=LGRAY)

    cd = ChartData()
    cd.categories = ["Synchronous (REST chain)", "Asynchronous (event-driven)"]
    cd.add_series("Client wait time (ms)", (720, 48))
    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(7.55), Inches(1.20), Inches(5.52), Inches(4.32), cd)
    chart = cf.chart
    chart.has_legend = False
    chart.has_title  = True
    chart.chart_title.text_frame.text = "Client wait time (ms)"
    try:
        chart.series[0].points[0].format.fill.solid()
        chart.series[0].points[0].format.fill.fore_color.rgb = RED
        chart.series[0].points[1].format.fill.solid()
        chart.series[0].points[1].format.fill.fore_color.rgb = GREEN
    except Exception:
        pass

    _rect(slide, 7.55, 5.64, 5.52, 0.48, fill=LT_GRN, stroke=GREEN, sw=Pt(1))
    _txt(slide, "720 ms  vs  48 ms  —  15x faster HTTP response to the client",
         7.75, 5.72, 5.12, 0.32, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "[~60 sec] "
        "The table shows the key trade-offs. Synchronous gives strong consistency and a "
        "simple call stack — but tight coupling and cascading failures. "
        "Async gives loose coupling and immediate responses — but eventual consistency "
        "and distributed debugging. The bar chart shows the difference: "
        "720 ms versus 48 ms to the client. "
        "The rule: choose async when the caller does NOT need the result right now."
    )


def slide_05_saga(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Choreography Saga Pattern",
            "Distributed transactions without a central orchestrator")
    _footer(slide, 5, total=10)

    # Happy path band
    _rect(slide, 0.30, 1.18, 6.22, 0.96, fill=LT_GRN, stroke=GREEN, sw=Pt(0.9))
    _txt(slide, "Happy path", 0.50, 1.24, 2.0, 0.26, size=11, bold=True, color=GREEN)

    BW, BH, HGAP = 1.40, 0.56, 0.22
    happy = [
        ("Order\nPlaced",     C_EVT,                    None),
        ("Stock\nReserved",   C_STK,                    None),
        ("Payment\nSucceeded",C_PAY,                    None),
        ("Order\nConfirmed",  RGBColor(0x0A,0x42,0x1A), None),
    ]
    hx = 0.46
    for i, (evt, bg, _) in enumerate(happy):
        _label(slide, evt, hx, 1.46, BW, BH, fill=bg, stroke=None, size=10, sw=Pt(0))
        if i < len(happy)-1:
            _arrow(slide, hx+BW, 1.74, hx+BW+HGAP, 1.74, color=GREEN, w=Pt(1.4))
        hx += BW + HGAP

    # Failure path band
    _rect(slide, 0.30, 2.30, 6.22, 1.20, fill=LT_RED, stroke=RED, sw=Pt(0.9))
    _txt(slide, "Failure:  PaymentFailed  ->  compensating transactions",
         0.50, 2.38, 5.90, 0.28, size=11, bold=True, color=RED)

    FBW, FBH, FGAP = 1.80, 0.56, 0.22
    fail = [
        ("Order\nPlaced",   C_EVT,                    None),
        ("Stock\nReserved", C_STK,                    None),
        ("Payment\nFailed", RGBColor(0x98,0x0A,0x0A), None),
    ]
    fx = 0.46
    pf_cx = 0.0
    for i, (evt, bg, _) in enumerate(fail):
        _label(slide, evt, fx, 2.70, FBW, FBH, fill=bg, stroke=None, size=10.5, sw=Pt(0))
        if i < len(fail)-1:
            _arrow(slide, fx+FBW, 2.98, fx+FBW+FGAP, 2.98, color=ACNT, w=Pt(1.4))
        if i == len(fail)-1:
            pf_cx = fx + FBW/2
        fx += FBW + FGAP

    pf_by = 2.70 + FBH   # 3.26

    # Compensating transactions (below the failure band)
    _txt(slide, "Compensating transactions:", 0.30, 3.64, 4.5, 0.26,
         size=10.5, bold=True, color=RED)

    comp = [
        (0.30, 3.94, 2.82, "Stock\nReleased\n(reservation undone)", C_STK),
        (3.40, 3.94, 2.82, "Order\nCancelled\n(order closed)",       C_API),
    ]
    for lx, ty, bw, name, col in comp:
        top_cx = lx + bw/2
        _arrow(slide, pf_cx, pf_by, top_cx, ty, color=RED, w=Pt(1.2))
        _label(slide, name, lx, ty, bw, 0.80, fill=col, stroke=RED, size=11, sw=Pt(1))

    # Code box (right)
    _txt(slide, "Compensation consumer (StockConsumer):",
         6.72, 1.18, 6.32, 0.28, size=11.5, bold=True)
    _code_box(slide, [
        "public async Task Consume(",
        "    ConsumeContext<PaymentFailed> context)",
        "{",
        "    var msg = context.Message;",
        "    // release the held inventory",
        "    await context.Publish(",
        "        new StockReleased {",
        "            OrderId       = msg.OrderId,",
        "            CorrelationId = msg.CorrelationId",
        "        });",
        "}",
    ], 6.72, 1.52, 6.32, 3.10, size=10.5)

    _rect(slide, 6.72, 4.74, 6.32, 2.40, fill=WHITE, stroke=MGRAY, sw=Pt(0.9))
    _txt(slide, "Choreography vs Orchestration",
         6.92, 4.84, 5.90, 0.30, size=12, bold=True)
    for i, f in enumerate([
        "No distributed lock — each service commits locally",
        "Compensation replaces rollback — forward-only",
        "Choreography (our): no central controller",
        "Orchestration alternative: saga manager directs steps",
    ]):
        _txt(slide, f"·  {f}", 6.92, 5.24+i*0.44, 5.90, 0.40, size=11, color=LGRAY)

    slide.notes_slide.notes_text_frame.text = (
        "[~80 sec] "
        "The Saga pattern replaces a distributed transaction with local commits. "
        "No 2-Phase Commit, no distributed locks. Each service commits locally and "
        "publishes the next event. Top row: the happy path. Bottom row: payment fails — "
        "StockConsumer hears PaymentFailed and publishes StockReleased, releasing the "
        "reserved inventory. OrderService publishes OrderCancelled. "
        "These are compensating transactions — they undo earlier work going forward, "
        "not backwards. We use choreography: no central controller; each service knows "
        "only its own responsibilities."
    )


def slide_06_idempotency(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Idempotency — Handling Duplicate Messages",
            "RabbitMQ guarantees at-least-once delivery, not exactly-once")
    _footer(slide, 6, total=10)

    _txt(slide, "Consumer decision logic:", 0.32, 1.22, 4.0, 0.28, size=12, bold=True)

    FC_L  = 0.52
    FC_W  = 3.06
    FC_CX = FC_L + FC_W/2   # 2.05"

    # Step 1: message arrives
    _label(slide, "Message arrives\n(with MessageId)",
           FC_L, 1.58, FC_W, 0.62, fill=C_API, stroke=None, size=11, sw=Pt(0))
    _arrow(slide, FC_CX, 2.20, FC_CX, 2.44, color=ACNT)

    # Diamond decision shape (type 4)
    dh = 1.10
    sh_d = slide.shapes.add_shape(4, Inches(FC_L), Inches(2.44), Inches(FC_W), Inches(dh))
    sh_d.fill.solid(); sh_d.fill.fore_color.rgb = LT_BLU
    sh_d.line.color.rgb = ACNT; sh_d.line.width = Pt(1.2)
    tf = sh_d.text_frame; tf.word_wrap = True
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Already in\nProcessedMessages?"
    r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = TEXT

    d_right  = FC_L + FC_W    # 3.58"
    d_bottom = 2.44 + dh      # 3.54"
    d_cy     = 2.44 + dh/2    # 2.99"

    # YES branch → right
    _arrow(slide, d_right, d_cy, 4.30, d_cy, color=RED, w=Pt(1.3))
    _txt(slide, "YES", d_right+0.06, d_cy-0.26, 0.62, 0.22, size=10, bold=True, color=RED)
    _label(slide, "Discard\n(return)", 4.30, d_cy-0.34, 1.88, 0.68,
           fill=LT_RED, stroke=RED, color=RED, size=11, bold=True, sw=Pt(1))

    # NO branch → down
    _arrow(slide, FC_CX, d_bottom, FC_CX, 3.74, color=GREEN, w=Pt(1.3))
    _txt(slide, "NO", FC_CX+0.08, d_bottom+0.04, 0.50, 0.22, size=10, bold=True, color=GREEN)

    _label(slide, "INSERT MessageId\n(unique constraint)",
           FC_L, 3.74, FC_W, 0.62, fill=C_STK, stroke=None, size=11, sw=Pt(0))
    _arrow(slide, FC_CX, 4.36, FC_CX, 4.56, color=GREEN)
    _label(slide, "Process message\n(business logic)",
           FC_L, 4.56, FC_W, 0.62, fill=C_STK, stroke=None, size=11, sw=Pt(0))

    # DB table illustration
    _rect(slide, 0.32, 5.36, 6.14, 1.14, fill=WHITE, stroke=MGRAY, sw=Pt(0.8))
    _txt(slide, "ProcessedMessages", 0.52, 5.44, 4.0, 0.28, size=10.5, bold=True, color=ACNT)
    _txt(slide, "MessageId (PK, UNIQUE)                   ProcessedAt",
         0.52, 5.74, 5.80, 0.24, size=9.5, bold=True)
    _line(slide, 0.52, 6.00, 6.38, 6.00, color=MGRAY, w=Pt(0.5))
    _txt(slide, "a9f1-4d22-...-8c3b                       2026-06-14 09:41",
         0.52, 6.02, 5.80, 0.22, size=9.5, color=LGRAY)
    _txt(slide, "a9f1-4d22-...-8c3b   DUPLICATE  -> DbUpdateException",
         0.52, 6.24, 5.80, 0.22, size=9.5, color=RED, bold=True)

    # Code box (right)
    _txt(slide, "StockConsumer — idempotency guard:",
         6.68, 1.22, 6.42, 0.28, size=12, bold=True)
    _code_box(slide, [
        "_db.ProcessedMessages.Add(",
        "    new ProcessedMessage {",
        "        MessageId   = context.MessageId",
        "                      ?? Guid.Empty,",
        "        ProcessedAt = DateTime.UtcNow",
        "    });",
        "try {",
        "    await _db.SaveChangesAsync();",
        "}",
        "catch (DbUpdateException) {",
        "    _logger.LogWarning(",
        '        "Duplicate {Id} - discarded",',
        "        context.MessageId);",
        "    return; // ack, don't process twice",
        "}",
        "// normal business logic...",
    ], 6.68, 1.58, 6.42, 5.68, size=10.5)

    slide.notes_slide.notes_text_frame.text = (
        "[~60 sec] "
        "RabbitMQ guarantees at-least-once delivery — not exactly-once. "
        "Every consumer must therefore be idempotent. Our mechanism: before any business "
        "logic, we try to INSERT the MessageId into a ProcessedMessages table that has a "
        "unique constraint. If the row already exists, we get DbUpdateException, log a "
        "warning, and return without processing — the message is acknowledged so it is "
        "not requeued. First delivery processes; every retry is silently discarded."
    )


def slide_07_retries(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Resilience — Retries & Dead-Letter Queues",
            "Transient failures are expected — handle them without losing messages")
    _footer(slide, 7, total=10)

    _txt(slide, "Retry timeline  (ShippingConsumer,  SHIPPING_FAILS=true):",
         0.32, 1.22, 11.0, 0.28, size=12, bold=True)

    # Timeline — boxes start at y=1.60, clear of subtitle (ends at 1.50")
    TL_Y  = 2.44   # baseline (box bottoms + dot line)
    BOX_H = 0.82
    BOX_W = 2.44

    attempts = [
        (0.32,  "Attempt 1\n(fails)",            RED,    LT_RED),
        (3.00,  "Wait 2 s\nAttempt 2\n(fails)",  RED,    LT_RED),
        (5.68,  "Wait 2 s\nAttempt 3\n(fails)",  RED,    LT_RED),
        (8.56,  "Dead-Letter\nQueue (_error)",    ORANGE, LT_ORG),
    ]
    _line(slide, 0.32, TL_Y, 12.80, TL_Y, color=MGRAY, w=Pt(1))
    for ax, label, stroke, fill in attempts:
        _label(slide, label, ax, TL_Y-BOX_H, BOX_W, BOX_H,
               fill=fill, stroke=stroke, color=stroke,
               size=11.5, bold=True, sw=Pt(1.2))
        cx = ax + BOX_W/2
        _rect(slide, cx-0.07, TL_Y-0.07, 0.14, 0.14, fill=stroke)

    # Arrows between attempt boxes
    gaps = [(a[0]+BOX_W, b[0]) for a, b in zip(attempts, attempts[1:])]
    for gx1, gx2 in gaps:
        _arrow(slide, gx1, TL_Y-BOX_H/2, gx2, TL_Y-BOX_H/2, color=RED, w=Pt(1.1))

    # Info bar
    _rect(slide, 0.32, TL_Y+0.10, 12.70, 0.44, fill=LT_ORG, stroke=ORANGE, sw=Pt(0.8))
    _txt(slide,
         "No message ever lost  ·  inspect in RabbitMQ UI at :15672  "
         "·  fix root cause  ·  requeue  ->  shipping-service_error",
         0.52, TL_Y+0.17, 12.30, 0.34, size=11, color=ORANGE)

    sec_y = TL_Y + 0.66
    _txt(slide, "MassTransit retry + DLQ (ShippingConsumer/Program.cs):",
         0.32, sec_y, 7.0, 0.28, size=11.5, bold=True)
    _code_box(slide, [
        'cfg.ReceiveEndpoint("shipping-service", e =>',
        "{",
        "    // 3 retries, 2 seconds apart",
        "    e.UseMessageRetry(r =>",
        "        r.Interval(3, TimeSpan.FromSeconds(2)));",
        "",
        "    // after all retries: moves to",
        "    // shipping-service_error (DLQ)",
        "    e.ConfigureConsumer",
        "        <OrderConfirmedConsumer>(context);",
        "});",
    ], 0.32, sec_y+0.34, 6.20, 3.30, size=11)

    _txt(slide, "RabbitMQ queues:", 6.72, sec_y, 6.30, 0.28, size=11.5, bold=True)
    queues = [
        ("shipping-service",       GREEN,  "active — normal processing"),
        ("shipping-service_error", ORANGE, "dead-lettered after 3 retries"),
        ("invoice-service",        GREEN,  "active"),
        ("warehouse-service",      GREEN,  "active"),
        ("notification-service",   GREEN,  "active"),
    ]
    qy = sec_y + 0.34
    for qname, qcol, qnote in queues:
        _rect(slide, 6.72, qy, 6.30, 0.52, fill=WHITE, stroke=qcol, sw=Pt(1.2))
        _txt(slide, qname, 6.90, qy+0.06, 3.60, 0.28, size=10.5, bold=True, color=qcol)
        _txt(slide, qnote, 6.90, qy+0.06, 5.90, 0.28, size=10,
             color=ORANGE if qcol == ORANGE else LGRAY, align=PP_ALIGN.RIGHT)
        qy += 0.60

    slide.notes_slide.notes_text_frame.text = (
        "[~55 sec] "
        "Transient failures — a network blip, a database timeout — are normal in "
        "distributed systems. Our mechanism: MassTransit retries up to three times "
        "with two seconds between attempts. If all three fail, the message moves to "
        "shipping-service_error — a dead-letter queue visible in the RabbitMQ UI on "
        "port 15672. An operator can inspect the message, fix the root cause, and "
        "requeue it. No message is ever silently lost."
    )


def slide_08_fanout(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Fan-Out — Publish-Subscribe",
            "One event triggers four independent consumers simultaneously")
    _footer(slide, 8, total=10)

    PUB_CY = 3.46   # vertical centre of publisher and event boxes

    # Publisher box (left)
    _label(slide, "Order Service\n(publisher)",
           0.40, PUB_CY-0.46, 2.20, 0.92,
           fill=C_API, stroke=None, size=12, sw=Pt(0))

    # Arrow publisher -> event  (x1=2.60 < x2=3.10 => goes RIGHT)
    _arrow(slide, 2.60, PUB_CY, 3.10, PUB_CY, color=ACNT, w=Pt(1.6))

    # Event box (centre)
    EV_L, EV_W, EV_H = 3.10, 2.68, 0.92
    _label(slide, "OrderConfirmed",
           EV_L, PUB_CY-EV_H/2, EV_W, EV_H,
           fill=RGBColor(0x0A,0x42,0x1A), stroke=None, size=15, bold=True, sw=Pt(0))
    _txt(slide, "RabbitMQ fanout exchange",
         EV_L, PUB_CY+EV_H/2+0.04, EV_W, 0.24,
         size=9, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

    EV_RIGHT = EV_L + EV_W   # 5.78"

    # 4 consumers — evenly spaced, all visible
    CL, CW, CH = 7.20, 5.92, 0.96
    c_tops = [1.28, 2.52, 3.76, 5.00]
    c_labels = [
        "Invoice Consumer\n(generates PDF invoice)",
        "Warehouse Consumer\n(picks items for the order)",
        "Shipping Consumer\n(creates shipping label)",
        "Notification Consumer\n(emails order confirmation)",
    ]
    for ty, name in zip(c_tops, c_labels):
        cy = ty + CH/2
        _arrow(slide, EV_RIGHT, PUB_CY, CL, cy, color=GREEN, w=Pt(1.5))
        _label(slide, name, CL, ty, CW, CH, fill=C_FAN, stroke=None, size=12, sw=Pt(0))

    # Key insight bar at bottom
    _rect(slide, 0.30, 6.24, 12.74, 0.82, fill=LT_BLU, stroke=ACNT, sw=Pt(1))
    _txt(slide,
         "Publisher does NOT know who is listening  —  adding a new consumer "
         "requires zero changes to Order Service.\n"
         "Each consumer has its own named queue  —  slow Invoice does NOT block Shipping.",
         0.50, 6.30, 12.34, 0.72, size=12, color=TEXT)

    slide.notes_slide.notes_text_frame.text = (
        "[~50 sec] "
        "OrderService publishes OrderConfirmed to a RabbitMQ fanout exchange. "
        "RabbitMQ delivers a copy to every bound queue — one per consumer. "
        "All four run in parallel, completely independently. "
        "The key property: the publisher does not know who is listening. "
        "Adding a new consumer means creating a new service and binding its queue — "
        "zero changes to OrderService. Slow Invoice processing does not block Shipping."
    )


def slide_09_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Testing Strategy",
            "Four tools, four layers — every pattern exercised")
    _footer(slide, 9, total=10)

    cells = [
        {"tool": "xUnit 2.9.2",        "role": "Test runner & assertions",
         "hdr": C_EVT, "stroke": ACNT,
         "code": ["[Fact]", "public async Task",
                  "  Consume_InStock_PublishesReserved()",
                  "{", "    var consumer =",
                  "      new OrderPlacedConsumer(", "        _logger, _db);",
                  "", "    await consumer.Consume(", "        _ctx.Object);",
                  "    // verified in Moq below", "}"]},
        {"tool": "Moq 4.20.72",         "role": "Mock dependencies & verify calls",
         "hdr": C_PAY, "stroke": RGBColor(0xA0,0x60,0xFF),
         "code": ["var ctx = new Mock<",
                  "  ConsumeContext<OrderPlaced>>();",
                  "ctx.Setup(x => x.Message)", "   .Returns(order);",
                  "", "ctx.Verify(x => x.Publish(",
                  "  It.Is<StockReserved>(r =>", "    r.OrderId == id),",
                  "  It.IsAny<CancellationToken>()),", "  Times.Once);"]},
        {"tool": "WireMock.Net 1.6.9",  "role": "Real HTTP stub for payment gateway",
         "hdr": RGBColor(0x76,0x3C,0x00), "stroke": ORANGE,
         "code": ["_server", "  .Given(Request.Create()",
                  '    .WithPath("/charge")', "    .UsingPost())",
                  "  .RespondWith(", "    Response.Create()",
                  "      .WithStatusCode(200));",
                  "", "var ok = await _client",
                  "  .ChargeAsync(id, 49.99m);", "Assert.True(ok);"]},
        {"tool": "Testcontainers 4.0.0","role": "Real PostgreSQL — constraint test",
         "hdr": C_STK, "stroke": GREEN,
         "code": ["private readonly",
                  "  PostgreSqlContainer _pg =",
                  "    new PostgreSqlBuilder()",
                  '      .WithDatabase("stockdb")',
                  "      .Build();",
                  "", "// second insert — same MessageId",
                  "await Assert.ThrowsAsync<",
                  "  DbUpdateException>(async () =>",
                  "    await db2", "      .SaveChangesAsync());"]},
    ]

    positions = [(0.30, 1.20), (6.82, 1.20), (0.30, 4.24), (6.82, 4.24)]
    CW, CH = 6.20, 2.92

    for cell, (lx, ty) in zip(cells, positions):
        _rect(slide, lx, ty, CW, CH, fill=WHITE, stroke=MGRAY, sw=Pt(1))
        _rect(slide, lx, ty, CW, 0.60, fill=cell["hdr"])
        _txt(slide, cell["tool"], lx+0.14, ty+0.07, CW-0.28, 0.36,
             size=13, bold=True, color=WHITE)
        _txt(slide, cell["role"], lx+0.14, ty+0.42, CW-0.28, 0.20,
             size=10, color=MGRAY, italic=True)
        _code_box(slide, cell["code"], lx+0.10, ty+0.68, CW-0.20, CH-0.78, size=9.2)

    slide.notes_slide.notes_text_frame.text = (
        "[~65 sec] "
        "Four tools, four layers. xUnit runs the tests. Moq mocks ConsumeContext "
        "so we test consumer logic without a live broker — we verify that Publish "
        "was called exactly once with the correct payload. WireMock.Net starts a "
        "real HTTP server on a random port to test PaymentGatewayClient — we stub "
        "POST /charge to return 200 or 402 and assert the return value. "
        "Testcontainers starts a real PostgreSQL 16 container to verify the unique "
        "constraint — an InMemory database cannot enforce that."
    )


def slide_10_docker(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Docker & Deployment",
            "docker compose up --build  ->  entire system running in one command")
    _footer(slide, 10, total=10)

    # Container grid — CH=0.82, CG=0.12 keeps rows clear of info cards below
    CW, FW = 2.92, 2.10
    CH, CG = 0.82, 0.12
    r1y = 1.22
    r2y = r1y + CH + CG   # 2.16
    r3y = r2y + CH + CG   # 3.10  → bottom = 3.92

    for i, (name, col) in enumerate([
        ("RabbitMQ 3\n:5672  ·  :15672", C_INF),
        ("PostgreSQL 16\n:5432",           C_INF),
        ("Stub Payment\nGateway  :9090",   RGBColor(0x5C,0x38,0x00)),
    ]):
        _label(slide, name, 0.35+i*(CW+CG), r1y, CW, CH,
               fill=col, stroke=None, size=11, sw=Pt(0))

    for i, (name, col) in enumerate([
        ("Order Service API\n:5000", C_API),
        ("Stock Consumer",            C_STK),
        ("Payment Consumer",          C_PAY),
    ]):
        _label(slide, name, 0.35+i*(CW+CG), r2y, CW, CH,
               fill=col, stroke=None, size=11, sw=Pt(0))

    for i, (name, col) in enumerate([
        ("Invoice\nConsumer",         C_FAN),
        ("Warehouse\nConsumer",       C_FAN),
        ("Shipping Consumer\n(+DLQ)", C_FAN),
        ("Notification\nConsumer",    C_FAN),
    ]):
        _label(slide, name, 0.35+i*(FW+CG), r3y, FW, CH,
               fill=col, stroke=None, size=10.5, sw=Pt(0))

    # Info cards — start at 4.10 (gap of 0.18" after container grid bottom 3.92)
    INFO_Y = 4.10
    _rect(slide, 0.30, INFO_Y, 5.96, 3.02, fill=WHITE, stroke=MGRAY, sw=Pt(0.8))
    _txt(slide, "Multi-stage Dockerfiles", 0.50, INFO_Y+0.10, 5.56, 0.30,
         size=12, bold=True, color=ACNT)
    _txt(slide,
         "Stage 1 (SDK image):\n"
         "  dotnet restore + dotnet publish\n\n"
         "Stage 2 (runtime image):\n"
         "  copy /app/publish only  ->  minimal image\n\n"
         "Build context = repo root\n"
         "  shared/Contracts/ accessible in all Dockerfiles",
         0.50, INFO_Y+0.46, 5.56, 2.48, size=11, color=LGRAY)

    _txt(slide, "Health check + dependency ordering:",
         6.46, INFO_Y, 6.68, 0.30, size=12, bold=True)
    _code_box(slide, [
        "rabbitmq:",
        "  healthcheck:",
        '    test: ["CMD",',
        '      "rabbitmq-diagnostics", "ping"]',
        "    interval: 10s",
        "    retries: 5",
        "",
        "stock-consumer:",
        "  depends_on:",
        "    rabbitmq:",
        "      condition: service_healthy",
        "    postgres:",
        "      condition: service_healthy",
    ], 6.46, INFO_Y+0.36, 6.68, 2.66, size=10.5)

    slide.notes_slide.notes_text_frame.text = (
        "[~50 sec] "
        "One command — docker compose up --build — starts all nine containers. "
        "Multi-stage Dockerfiles keep images small: SDK image for build, runtime image "
        "for execution. Health checks on RabbitMQ and PostgreSQL prevent race conditions "
        "where a consumer starts before the broker is ready. Configuration is injected "
        "via environment variables, so the same images run locally and in production."
    )


def slide_11_ux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "What Does the User Experience?",
            "Eventual consistency from the customer's perspective — 4 scenarios")
    _footer(slide, 11, total=10)

    scenarios = [
        {"label": "1  Happy Path", "col": GREEN, "bg": LT_GRN,
         "steps": [("POST /orders", WHITE,  C_API,  ACNT),
                   ("HTTP 201\nPending",  C_API,  WHITE,  ACNT),
                   ("Stock +\nPayment OK",C_STK,  WHITE,  GREEN),
                   ("Invoice · WH\nShip", C_FAN,  WHITE,  GREEN),
                   ("[email]\nConfirmed", LT_GRN, GREEN,  GREEN)]},
        {"label": "2  No Stock", "col": RED, "bg": LT_RED,
         "steps": [("POST /orders",       WHITE,  C_API,  ACNT),
                   ("HTTP 201\nPending",   C_API,  WHITE,  ACNT),
                   ("StockUnavailable\npubd",LT_RED,RED,   RED),
                   ("Payment\nnever called",DBOX, LGRAY,   MGRAY),
                   ("[email]\nCancelled",  LT_RED, RED,    RED)]},
        {"label": "3  Payment Fails  (saga compensation)", "col": ORANGE, "bg": LT_ORG,
         "steps": [("POST /orders",       WHITE,  C_API,   ACNT),
                   ("HTTP 201\nPending",   C_API,  WHITE,   ACNT),
                   ("Stock\nReserved",     C_STK,  WHITE,   GREEN),
                   ("PaymentFailed\nStockReleased",LT_ORG,ORANGE,ORANGE),
                   ("[email]\nCancelled",  LT_ORG, ORANGE,  ORANGE)]},
        {"label": "4  Shipping Fails  (retries + DLQ, customer unaffected)", "col": ACNT, "bg": LT_BLU,
         "steps": [("POST /orders",       WHITE,  C_API,  ACNT),
                   ("HTTP 201\nPending",   C_API,  WHITE,  ACNT),
                   ("Confirmed\nInvoice+WH",C_FAN, WHITE,  GREEN),
                   ("3x retry\n+ DLQ",     LT_ORG, ORANGE, ORANGE),
                   ("[email]\nConfirmed",  LT_GRN, GREEN,  GREEN)]},
    ]

    ROW_H = 1.20
    COL_W = 2.34
    SH    = 0.82
    for ri, sc in enumerate(scenarios):
        bt = 1.22 + ri*ROW_H
        _rect(slide, 0.30, bt, 0.08, SH+0.22, fill=sc["col"])
        _txt(slide, sc["label"], 0.46, bt, 12.50, 0.22,
             size=10, bold=True, color=sc["col"])
        for ci, (stxt, bg, tc, stroke) in enumerate(sc["steps"]):
            lx = 0.46 + ci*(COL_W+0.16)
            _label(slide, stxt, lx, bt+0.26, COL_W, SH-0.26,
                   fill=bg, stroke=stroke, color=tc,
                   size=9.5, bold=(bg not in (WHITE, DBOX)), sw=Pt(0.8))
            if ci < 4:
                _arrow(slide, lx+COL_W, bt+0.26+(SH-0.26)/2,
                       lx+COL_W+0.16, bt+0.26+(SH-0.26)/2,
                       color=sc["col"], w=Pt(1))

    _txt(slide,
         "User always gets immediate HTTP 201  ·  "
         "CorrelationId links all events across the entire flow",
         0.30, 7.10, 12.74, 0.26, size=11, color=LGRAY,
         align=PP_ALIGN.CENTER, italic=True)

    slide.notes_slide.notes_text_frame.text = (
        "[~50 sec] "
        "In every scenario the customer gets HTTP 201 immediately — status Pending. "
        "Happy path: confirmed within seconds. No stock: payment never called, order "
        "cancelled. Payment failure: saga compensation releases the reserved stock. "
        "Shipping failure: the customer still sees a confirmation email; the DLQ "
        "handles the failure silently — the operator fixes and requeues."
    )


def slide_12_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Key Takeaways",
            "Event-driven: loose coupling, resilience, scalability — with real trade-offs")
    _footer(slide, 10, total=10)

    # Metric row — 6 equal boxes centred
    metrics = [("7","microservices"),("6","demo scenarios"),("10","tests"),
               ("4","patterns"),("9","containers"),("0","shared DBs")]
    mw, mg = 1.92, 0.20
    ml = (13.33 - (len(metrics)*mw + (len(metrics)-1)*mg)) / 2
    for i, (num, lbl) in enumerate(metrics):
        lx = ml + i*(mw+mg)
        _rect(slide, lx, 1.22, mw, 1.12, fill=WHITE, stroke=ACNT, sw=Pt(1.5))
        _txt(slide, num, lx, 1.28, mw, 0.70,
             size=38, bold=True, color=ACNT, align=PP_ALIGN.CENTER)
        _txt(slide, lbl, lx, 1.98, mw, 0.28,
             size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    # Pattern badges — 4 equal boxes centred
    patterns = [
        ("Choreography Saga",   "Distributed consistency\nwithout 2PC",    C_PAY),
        ("Idempotent Consumer", "Safe at-least-once\ndelivery",             C_STK),
        ("Dead-Letter Queue",   "No message ever\nsilently lost",           C_INF),
        ("Pub/Sub Fan-Out",     "Independent parallel\nprocessing",         C_API),
    ]
    pw, pg = 2.90, 0.38
    pl = (13.33 - (len(patterns)*pw + (len(patterns)-1)*pg)) / 2
    for i, (name, desc, col) in enumerate(patterns):
        lx = pl + i*(pw+pg)
        _label(slide, name, lx, 2.62, pw, 0.56, fill=col, stroke=None, size=13, sw=Pt(0))
        _txt(slide, desc, lx+0.08, 3.24, pw-0.16, 0.56,
             size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    # Benefits / Trade-offs cards
    _rect(slide, 0.30, 4.02, 5.98, 2.18, fill=LT_GRN, stroke=GREEN, sw=Pt(1))
    _txt(slide, "Benefits", 0.50, 4.10, 5.58, 0.32, size=13, bold=True, color=GREEN)
    for i, p in enumerate([
        "Loose coupling — services are independent",
        "Resilience — failures are isolated",
        "Scalability — scale each service independently",
        "Extensibility — add consumers, publisher unchanged",
    ]):
        _txt(slide, f"·  {p}", 0.50, 4.50+i*0.38, 5.58, 0.34, size=11, color=GREEN)

    _rect(slide, 6.56, 4.02, 6.48, 2.18, fill=LT_RED, stroke=RED, sw=Pt(1))
    _txt(slide, "Trade-offs", 6.76, 4.10, 6.08, 0.32, size=13, bold=True, color=RED)
    for i, c in enumerate([
        "Eventual consistency — not strong",
        "Distributed debugging — 7 log streams",
        "Operational complexity — 9 containers",
        "Idempotency required — duplicates happen",
    ]):
        _txt(slide, f"·  {c}", 6.76, 4.50+i*0.38, 6.08, 0.34, size=11, color=RED)

    _rect(slide, 0.30, 6.36, 12.74, 0.76, fill=HDR)
    _txt(slide, "Thank you  —  I am happy to take questions",
         0.30, 6.44, 12.74, 0.60,
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "[~30 sec] "
        "To summarise: four patterns, seven services, ten tests, nine containers. "
        "Event-driven gives loose coupling, resilience, and independent scalability. "
        "The trade-offs are real: eventual consistency, distributed debugging, "
        "and mandatory idempotency. I am happy to discuss any of this."
    )


# ── Build ─────────────────────────────────────────────────────────────────────

def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for fn in [slide_01_title, slide_02_problem, slide_03_architecture,
               slide_04_sync_async, slide_05_saga, slide_06_idempotency,
               slide_07_retries, slide_08_fanout, slide_09_testing,
               slide_10_docker, slide_12_takeaways]:
        fn(prs)
    out = os.path.join(OUT, "presentation.pptx")
    prs.save(out)
    print(f"  PPTX  {out}  ({len(prs.slides)} slides)")


# ── PDF talking paper ─────────────────────────────────────────────────────────

NOTES = [
    # (title, spoken text ~60 sec target, key bullet points)
    ("Event-Driven Order Fulfillment System",
     "Good morning. My name is Marco and I will present my synopsis project: "
     "an Event-Driven Order Fulfillment System in .NET 10. "
     "I will walk through the architecture, four design patterns, and the testing strategy "
     "in about ten minutes.",
     ["Target: ~20 sec"]),

    ("Why Event-Driven?",
     "Synchronously, every service call blocks the next — the customer waits for the "
     "sum of all latencies, and if Shipping goes down the entire chain fails. "
     "Event-driven: the Order API publishes one event and returns HTTP 201 immediately. "
     "RabbitMQ routes a copy to every subscriber in parallel. "
     "Services are decoupled in time and space — one failure no longer cascades.",
     ["Target: ~60 sec",
      "POST /orders -> publishes OrderPlaced -> HTTP 201 immediately",
      "One failure does NOT cascade"]),

    ("System Architecture",
     "Seven microservices, one shared Contracts library, RabbitMQ as the broker, "
     "and PostgreSQL for state. The horizontal chain is the core saga: OrderService "
     "publishes OrderPlaced, StockConsumer reserves stock, PaymentConsumer charges the "
     "card, then OrderService publishes OrderConfirmed. That single event fans out via a "
     "RabbitMQ fanout exchange to four consumers in parallel — Invoice, Warehouse, "
     "Shipping, Notification — each with its own named queue. "
     "Failure path: if payment fails, compensating events undo earlier steps.",
     ["Target: ~75 sec",
      "7 services + RabbitMQ + PostgreSQL + StubPaymentGateway",
      "Shared Contracts library — compile-time type safety",
      "Failure path: PaymentFailed -> StockReleased + OrderCancelled"]),

    ("Synchronous vs. Asynchronous",
     "The table shows the key trade-offs. Synchronous gives strong consistency and a "
     "simple call stack — but tight coupling and cascading failures. "
     "Async gives loose coupling and immediate responses — but eventual consistency "
     "and distributed debugging. The bar chart: 720 ms versus 48 ms to the client. "
     "The rule: choose async when the caller does NOT need the result right now.",
     ["Target: ~60 sec",
      "Sync: tight coupling, cascading failure, 720 ms wait",
      "Async: loose coupling, independent failure, 48 ms ACK",
      "Rule: async when result not needed immediately"]),

    ("Choreography Saga Pattern",
     "The Saga pattern replaces a distributed transaction with local commits — no "
     "2-Phase Commit, no distributed locks. Each service commits locally and publishes "
     "the next event. Top row: the happy path. Bottom row: payment fails — "
     "StockConsumer hears PaymentFailed and publishes StockReleased, releasing the "
     "reserved inventory. OrderService publishes OrderCancelled. "
     "These are compensating transactions — forward-only, not rollback. "
     "We use choreography: no central controller; each service knows only its own role.",
     ["Target: ~80 sec",
      "No 2PC — each service commits locally",
      "PaymentFailed -> StockReleased + OrderCancelled (compensation)",
      "Choreography: no central controller"]),

    ("Idempotency — Duplicate Messages",
     "RabbitMQ guarantees at-least-once delivery — not exactly-once. "
     "Every consumer must be idempotent. Our mechanism: before any business logic, "
     "we INSERT the MessageId into a ProcessedMessages table with a unique constraint. "
     "If the row already exists, DbUpdateException is thrown — we log a warning and "
     "return without processing. The message is acknowledged so it is not requeued. "
     "First delivery processes; every retry is silently discarded.",
     ["Target: ~60 sec",
      "Unique constraint on ProcessedMessages.MessageId",
      "DbUpdateException -> catch -> return (ack, not requeued)"]),

    ("Resilience — Retries & Dead-Letter Queues",
     "Transient failures — a network blip, a database timeout — are normal. "
     "Our mechanism: MassTransit retries up to three times with two seconds between "
     "attempts. If all three fail, the message moves to shipping-service_error — a "
     "dead-letter queue visible in the RabbitMQ UI on port 15672. An operator can "
     "inspect the message, fix the root cause, and requeue it. "
     "No message is ever silently lost.",
     ["Target: ~55 sec",
      "e.UseMessageRetry(r => r.Interval(3, TimeSpan.FromSeconds(2)))",
      "DLQ: shipping-service_error — visible in RabbitMQ UI at :15672"]),

    ("Fan-Out — Publish-Subscribe",
     "OrderService publishes OrderConfirmed to a RabbitMQ fanout exchange. "
     "RabbitMQ delivers a copy to every bound queue — one per consumer. "
     "All four run in parallel, completely independently. "
     "The key property: the publisher does not know who is listening. "
     "Adding a new consumer means creating a new service and binding its queue — "
     "zero changes to OrderService. Slow Invoice processing does not block Shipping.",
     ["Target: ~50 sec",
      "Fanout exchange -> one message, N independent queues",
      "Adding a subscriber = zero changes to the publisher"]),

    ("Testing Strategy",
     "Four tools, four layers. xUnit runs the tests. Moq mocks ConsumeContext "
     "so we test consumer logic without a live broker — we verify that Publish was "
     "called exactly once with the correct payload. WireMock.Net starts a real HTTP "
     "server to test PaymentGatewayClient — we stub POST /charge to return 200 or 402. "
     "Testcontainers starts a real PostgreSQL 16 container to verify the unique "
     "constraint — an InMemory database cannot enforce that constraint.",
     ["Target: ~65 sec",
      "xUnit: [Fact] tests + assertions",
      "Moq: mock ConsumeContext<T>, verify Publish() Times.Once",
      "WireMock: real HTTP stub POST /charge -> 200 or 402",
      "Testcontainers: real Postgres, unique constraint verification"]),

    ("Docker & Deployment",
     "One command — docker compose up --build — starts all nine containers. "
     "Multi-stage Dockerfiles keep images small: SDK image for build, runtime image "
     "for execution. Health checks on RabbitMQ and PostgreSQL prevent race conditions "
     "where a consumer starts before the broker is ready. Configuration is injected "
     "via environment variables, so the same images run locally and in production.",
     ["Target: ~50 sec",
      "docker compose up --build — one command, nine containers",
      "Multi-stage Dockerfile: SDK build stage + runtime stage",
      "depends_on: condition: service_healthy"]),

    ("Key Takeaways",
     "To summarise: four patterns, seven services, ten tests, nine containers. "
     "Event-driven gives loose coupling, resilience, and independent scalability. "
     "The trade-offs are real: eventual consistency, distributed debugging, "
     "and mandatory idempotency. I am happy to discuss any of this.",
     ["Target: ~30 sec",
      "Benefits: loose coupling, resilience, independent scalability",
      "Trade-offs: eventual consistency, distributed debugging, idempotency required"]),
]

APPENDIX = [
    ("EasyNetQ vs MassTransit",
     "EasyNetQ is a simple RabbitMQ client — quick to start but requires manual retry and "
     "DLQ implementation. MassTransit provides built-in abstractions for all of these."),
    ("API Gateway (Ocelot)",
     "An API gateway provides a single entry point. Ocelot is configured via JSON — "
     "routing, authentication, rate limiting, and response aggregation."),
    ("Security — HashiCorp Vault",
     "Instead of hardcoding credentials in environment variables, a production system "
     "fetches secrets from Vault at startup. Vault provides dynamic secrets and audit logs."),
    ("Event Sourcing (differs from event-driven)",
     "This project is event-driven but NOT event-sourced. Event Sourcing stores every "
     "state change as an immutable event; current state is derived by replaying events."),
    ("Sidecar / BFF / Anti-Corruption Layer",
     "Sidecar: helper container alongside the main service. "
     "BFF: dedicated API layer per client type. "
     "ACL: translates between bounded contexts with different domain models."),
    ("Sam Newman — Building Microservices",
     "Key themes: database-per-service (we follow this), bounded contexts, "
     "and the danger of distributed monoliths. Our system communicates only through events."),
]


def build_pdf():
    out = os.path.join(OUT, "talking_paper.pdf")
    doc = SimpleDocTemplate(out, pagesize=A4,
                            leftMargin=2.5*cm, rightMargin=2.5*cm,
                            topMargin=2.5*cm, bottomMargin=2.5*cm)
    NAVY = colors.HexColor("#10203C")
    TEAL = colors.HexColor("#0088AA")
    MID  = colors.HexColor("#445566")
    LGRY = colors.HexColor("#F4F7FB")
    ss   = getSampleStyleSheet()
    mk   = lambda n, **kw: ParagraphStyle(n, parent=ss["Normal"], **kw)

    s_title = mk("T",  fontSize=22, leading=28, textColor=NAVY,
                 spaceAfter=4, alignment=TA_CENTER, fontName="Helvetica-Bold")
    s_sub   = mk("S",  fontSize=13, leading=18, textColor=TEAL,
                 spaceAfter=4, alignment=TA_CENTER, fontName="Helvetica")
    s_snum  = mk("SN", fontSize=10, leading=14, textColor=TEAL,
                 spaceAfter=2, fontName="Helvetica-Bold")
    s_head  = mk("SH", fontSize=15, leading=20, textColor=colors.white,
                 spaceAfter=4, spaceBefore=6, fontName="Helvetica-Bold",
                 backColor=NAVY, borderPadding=(6,8,6,8))
    s_bul   = mk("BU", fontSize=10, leading=14, textColor=MID, spaceAfter=2,
                 leftIndent=0.4*cm, fontName="Courier",
                 backColor=LGRY, borderPadding=(4,6,4,6))
    s_body  = mk("BO", fontSize=11, leading=17,
                 textColor=colors.HexColor("#1A2A3A"),
                 spaceAfter=8, alignment=TA_JUSTIFY, fontName="Helvetica")
    s_lbl   = mk("LA", fontSize=9,  leading=12, textColor=TEAL,
                 spaceAfter=2, fontName="Helvetica-Bold")
    s_intro = mk("IN", fontSize=11, leading=17, textColor=MID,
                 spaceAfter=10, alignment=TA_JUSTIFY, fontName="Helvetica-Oblique")
    s_ahead = mk("AH", fontSize=13, leading=17, textColor=NAVY,
                 spaceBefore=12, spaceAfter=4, fontName="Helvetica-Bold")

    # Time budget data — (slide_label, short_title, target_sec)
    _time_rows = [
        ("#",  "Slide",                          "Target", "Running"),
        ("1",  "Title",                          "0:20",   "0:20"),
        ("2",  "Why Event-Driven?",              "1:00",   "1:20"),
        ("3",  "System Architecture",            "1:15",   "2:35"),
        ("4",  "Sync vs. Async",                 "1:00",   "3:35"),
        ("5",  "Choreography Saga",              "1:20",   "4:55"),
        ("6",  "Idempotency",                    "1:00",   "5:55"),
        ("7",  "Retries & DLQ",                  "0:55",   "6:50"),
        ("8",  "Fan-Out / Pub-Sub",              "0:50",   "7:40"),
        ("9",  "Testing Strategy",               "1:05",   "8:45"),
        ("10", "Docker & Deployment",            "0:50",   "9:35"),
        ("11", "Key Takeaways",                  "0:30",  "10:05"),
    ]
    _col_w = [1.2*cm, 8.0*cm, 2.2*cm, 2.2*cm]
    tbl = Table(_time_rows, colWidths=_col_w)
    tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0, 0), (-1, 0),  colors.HexColor("#10203C")),
        ("TEXTCOLOR",   (0, 0), (-1, 0),  colors.white),
        ("FONTNAME",    (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTSIZE",    (0, 0), (-1, 0),  9),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1),
         [colors.HexColor("#F4F7FB"), colors.white]),
        ("FONTNAME",    (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE",    (0, 1), (-1, -1), 9),
        ("ALIGN",       (0, 0), (0, -1),  "CENTER"),
        ("ALIGN",       (2, 0), (-1, -1), "RIGHT"),
        ("GRID",        (0, 0), (-1, -1), 0.5, colors.HexColor("#C4D2E0")),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING",(0, 0), (-1, -1), 6),
        ("TOPPADDING",  (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING",(0,0), (-1, -1), 4),
    ]))

    story = [
        Spacer(1, 1.5*cm),
        Paragraph("Talking Paper", s_title),
        Paragraph("Event-Driven Order Fulfillment System", s_sub),
        Paragraph("System Integration · Spring 2026", s_sub),
        Spacer(1, 0.5*cm),
        HRFlowable(width="100%", thickness=2, color=TEAL),
        Spacer(1, 0.4*cm),
        Paragraph(
            "Full spoken script for each of the 11 slides. "
            "Exam format: individual external exam — presentation ~10 min, "
            "discussion ~15 min, total 30 min (Henrik Kuhl, SYS Spring 2026).",
            s_intro),
        Spacer(1, 0.3*cm),
        Paragraph("Time Budget", mk("TB", fontSize=11, leading=14, textColor=NAVY,
                                    spaceAfter=4, fontName="Helvetica-Bold")),
        tbl,
        Spacer(1, 0.4*cm),
        HRFlowable(width="100%", thickness=1, color=colors.lightgrey),
        PageBreak(),
    ]

    for i, (title, notes, bullets) in enumerate(NOTES):
        block = [Paragraph(f"SLIDE {i+1} OF {len(NOTES)}", s_snum),
                 Paragraph(title, s_head), Spacer(1, 0.2*cm)]
        if bullets:
            block.append(Paragraph("Key points:", s_lbl))
            for b in bullets:
                block.append(Paragraph(b, s_bul))
            block.append(Spacer(1, 0.25*cm))
        block.append(Paragraph("What to say:", s_lbl))
        for para in notes.strip().split("\n"):
            if para.strip():
                block.append(Paragraph(para, s_body))
        block += [Spacer(1, 0.2*cm),
                  HRFlowable(width="100%", thickness=0.5, color=colors.lightgrey)]
        story.extend(block)
        if i < len(NOTES) - 1:
            story.append(PageBreak())

    story += [PageBreak(), Paragraph("Appendix: Curriculum Topics", s_title),
              HRFlowable(width="100%", thickness=2, color=TEAL), Spacer(1, 0.4*cm)]
    for heading, text in APPENDIX:
        story.append(KeepTogether([
            Paragraph(heading, s_ahead),
            Paragraph(text, s_body),
        ]))

    doc.build(story)
    print(f"  PDF   {out}")


if __name__ == "__main__":
    print("Generating...")
    build_pptx()
    build_pdf()
    print("Done.")
